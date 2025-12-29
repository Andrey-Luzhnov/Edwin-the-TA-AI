from flask import json
import fitz  # PyMuPDF (for PDF text/images)
from pptx import Presentation
import io
import base64
import uuid

# Snowflake Cortex Configuration
CORTEX_MODEL = "llama3-70b"  # Options: llama3-70b, llama3-8b, mistral-large, mixtral-8x7b

# -----------------------------
# Helper: Build baseline context
# -----------------------------
def get_baseline_context(courseID, connection, limit_chars=800):
    """
    Fetch course materials from DB and build system baseline context.
    OPTIMIZED: Reduced to 800 chars to avoid Snowflake 8192 token limit.
    """
    cursor = connection.cursor()
    # OPTIMIZATION: Only fetch 3 most recent materials
    cursor.execute(
        "SELECT title, content FROM course_materials WHERE course_id = %s ORDER BY created_at DESC LIMIT 3",
        (courseID,)
    )
    materials = cursor.fetchall()
    cursor.close()

    baseline = f"""You are Edwin, TA AI for course {courseID}. Give helpful answers. Reference course materials when relevant.

COURSE MATERIALS:
"""

    for m in materials:
        title, content = m
        if content:
            baseline += f"\n{title}: {content[:limit_chars]}\n"

    return baseline


def ingest_pdf_to_snowflake(file_path, courseID, connection):
    """Extract text from PDF and insert as ONE material entry."""
    import os

    doc = fitz.open(file_path)
    full_text = ""
    cursor = connection.cursor()

    # Extract text from all pages
    for page_num, page in enumerate(doc):
        text = page.get_text()
        full_text += f"--- Page {page_num + 1} ---\n{text}\n\n"

    # Get the PDF filename (without path)
    pdf_filename = os.path.basename(file_path)

    # Store as ONE material entry with the actual filename
    cursor.execute(
        "INSERT INTO course_materials (course_id, title, content) VALUES (%s, %s, %s)",
        (courseID, pdf_filename, full_text)
    )

    connection.commit()
    cursor.close()
    return True, f"PDF '{pdf_filename}' ingested successfully as one material", None


def ingest_pptx_to_snowflake(file_path, courseID, connection):
    """Extract text from PPTX and insert into course_materials."""
    prs = Presentation(file_path)
    cursor = connection.cursor()

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slide_text = "\n".join(texts)
        print("Ready#")
        cursor.execute(
            "INSERT INTO course_materials (course_id, title, content) VALUES (%s, %s, %s)",
            (courseID, f"Slide {i+1}", slide_text)
        )

    connection.commit()
    cursor.close()
    return True, "PPTX ingested successfully", None


# ---------------------------------------
# Create a blank conversation (no user yet)
# ---------------------------------------
def create_blank_conversation(courseID, connection, model=None):
    """
    Create a blank conversation with preloaded baseline course materials.
    Stored in DB without a user_id (unassigned).
    With Cortex, we just generate a unique conversation ID and store context.
    """
    # Generate unique conversation ID
    conv_id = f"conv_{uuid.uuid4().hex}"
    print("NEW BLANK CONVERSATION:", conv_id)

    # Get baseline context (course materials)
    baseline_context = get_baseline_context(courseID, connection)

    # Store conversation in DB with baseline context
    cursor = connection.cursor()
    cursor.execute(
        "INSERT INTO conversations (course_id, user_id, conv_id, is_assigned) VALUES (%s, NULL, %s, FALSE)",
        (courseID, conv_id)
    )
    connection.commit()

    # Store baseline context as first system message
    cursor.execute(
        "INSERT INTO edwin_messages (conv_id, user_id, userorAI, message) VALUES (%s, NULL, TRUE, %s)",
        (conv_id, baseline_context)
    )
    connection.commit()
    cursor.close()

    return conv_id


# ------------------------------------
# Start a new thread (assign to a user)
# ------------------------------------
def start_new_thread(user_id, courseID, connection):
    """
    Assigns a user to the next available blank conversation for this course.
    """
    cursor = connection.cursor()

    # find the first unassigned conversation
    cursor.execute("""
        SELECT id, conv_id FROM conversations
        WHERE course_id = %s AND is_assigned = FALSE
        LIMIT 1
    """, (courseID,))
    row = cursor.fetchone()

    if not row:
        cursor.close()
        return False, "No blank conversations available. Pre-generate more.", None

    conv_id = row[1]  # conv_id

    # assign to user
    cursor.execute("""
        UPDATE conversations
        SET user_id = %s, is_assigned = TRUE
        WHERE id = %s
    """, (user_id, row[0]))
    connection.commit()
    cursor.close()

    return True, "New thread started", conv_id


# -------------------------------
# Get the user's active conversation
# -------------------------------
def get_user_conversation(user_id, courseID, connection):
    cursor = connection.cursor()
    cursor.execute("""
        SELECT conv_id FROM conversations
        WHERE course_id = %s AND user_id = %s AND is_assigned = TRUE
        ORDER BY created_at DESC LIMIT 1
    """, (courseID, user_id))
    row = cursor.fetchone()
    cursor.close()
    if row:
        return row[0]
    return None


# -------------------------------
# SNOWFLAKE logging
# -------------------------------
def log_to_snowflake(conv_id, user_id, userorai, message, connection):
    # IF USER SENT MESSAGE, userorai is False. Else True for Edwin.
    cursor = connection.cursor()
    cursor.execute(
        "INSERT INTO edwin_messages (conv_id, user_id, userorAI, message) VALUES (%s, %s, %s, %s)",
        (conv_id, user_id, userorai, message)
    )
    connection.commit()
    cursor.close()


# -------------------------------
# Get conversation history
# -------------------------------
def get_conversation_history(conv_id, connection, limit=3):
    """
    Retrieve the last N messages from a conversation.
    Returns formatted string for context.
    OPTIMIZED: Reduced to 3 messages to avoid Snowflake 8192 token limit.
    """
    cursor = connection.cursor()
    cursor.execute("""
        SELECT userorAI, message, created_at
        FROM edwin_messages
        WHERE conv_id = %s
        ORDER BY created_at DESC
        LIMIT %s
    """, (conv_id, limit))

    messages = cursor.fetchall()
    cursor.close()

    # Reverse to get chronological order
    messages = list(reversed(messages))

    # Format conversation history (OPTIMIZATION: Heavily limit message length)
    history = ""
    for is_ai, message, timestamp in messages:
        if is_ai:
            # Skip the first system message (baseline context) from history
            if "You are Edwin" in message or "TA AI" in message:
                continue
            # OPTIMIZATION: Truncate to 200 chars
            history += f"Edwin: {message[:200]}\n"
        else:
            # OPTIMIZATION: Truncate to 150 chars
            history += f"Student: {message[:150]}\n"

    return history


# --------------------------
# Retrieve relevant course materials (RAG)
# --------------------------
def retrieve_relevant_materials(courseID, question, connection, limit=2):
    """
    Retrieve most relevant course materials for the question.
    Returns list of (title, content_snippet, source_url, full_content)
    OPTIMIZED: Reduced to 2 materials max, 500 chars each to avoid token limit.
    """
    cursor = connection.cursor()

    # OPTIMIZATION: Only fetch 10 most recent materials
    cursor.execute("""
        SELECT material_id, title, content, source_url
        FROM course_materials
        WHERE course_id = %s AND content IS NOT NULL
        ORDER BY created_at DESC
        LIMIT 10
    """, (courseID,))

    materials = cursor.fetchall()
    cursor.close()

    if not materials:
        return []

    # Simple keyword matching for relevance scoring
    question_lower = question.lower()
    # OPTIMIZATION: Filter out common words for better matching
    stop_words = {'the', 'is', 'at', 'which', 'on', 'a', 'an', 'as', 'are', 'was', 'were', 'be', 'been', 'what', 'when', 'where', 'who', 'how'}
    question_words = set(w for w in question_lower.split() if w not in stop_words and len(w) > 2)

    scored_materials = []
    for material_id, title, content, source_url in materials:
        if not content:
            continue

        content_lower = content.lower()
        title_lower = title.lower()

        # OPTIMIZATION: Score based on keyword overlap
        score = 0

        # Boost score if question words appear in title (highest priority)
        title_boost = sum(10 for word in question_words if word in title_lower)
        score += title_boost

        # Quick content check - count matches in first 500 chars only
        content_sample = content_lower[:500]
        overlap = sum(1 for word in question_words if word in content_sample)
        score += overlap

        # OPTIMIZATION: Only process materials with minimum score
        if score > 0:
            # Extract relevant snippet (smaller for speed)
            snippet = content[:150]
            for word in question_words:
                if word in content_lower:
                    idx = content_lower.find(word)
                    start = max(0, idx - 75)
                    end = min(len(content), idx + 75)
                    snippet = "..." + content[start:end] + "..."
                    break

            scored_materials.append({
                'score': score,
                'title': title,
                'snippet': snippet,
                'source_url': source_url or '',
                'full_content': content[:500]  # OPTIMIZATION: Reduced from 1500 to 500 chars
            })

    # Sort by score and return top N
    scored_materials.sort(key=lambda x: x['score'], reverse=True)
    return scored_materials[:limit]


# --------------------------
# Ask a question in a thread using Snowflake Cortex
# --------------------------
def ask_question(user_id, courseID, question, connection, model=None):
    """
    Ask a question using Snowflake Cortex AI with RAG citations.
    Returns: (success, message, answer_or_dict)
    answer_or_dict format: {"answer": "...", "citations": [...]}
    """
    conv_id = get_user_conversation(user_id, courseID, connection)
    if not conv_id:
        return False, "No active thread found. Start a new one first.", None

    # Log user question
    log_to_snowflake(conv_id, user_id, False, question, connection)

    # Get conversation history
    history = get_conversation_history(conv_id, connection)

    # Retrieve relevant course materials (RAG)
    relevant_materials = retrieve_relevant_materials(courseID, question, connection, limit=3)

    # Build context from relevant materials
    materials_context = ""
    citations = []

    if relevant_materials:
        materials_context = "\n\nRELEVANT COURSE MATERIALS:\n"
        for i, mat in enumerate(relevant_materials, 1):
            materials_context += f"\n[Source {i}] {mat['title']}:\n{mat['full_content']}\n"
            citations.append({
                'title': mat['title'],
                'url': mat['source_url'],
                'snippet': mat['snippet']
            })

    # Get baseline context (first message in conversation)
    cursor = connection.cursor()
    cursor.execute("""
        SELECT message FROM edwin_messages
        WHERE conv_id = %s AND userorAI = TRUE
        ORDER BY created_at ASC
        LIMIT 1
    """, (conv_id,))
    baseline_row = cursor.fetchone()
    baseline_context = baseline_row[0] if baseline_row else ""

    # Build full prompt for Cortex with RAG context
    full_prompt = f"""{baseline_context}

{materials_context}

CONVERSATION HISTORY:
{history}

Student: {question}

Edwin (cite sources when using information from course materials):"""

    # Use Snowflake Cortex to generate response
    cortex_model = model or CORTEX_MODEL
    cursor.execute("""
        SELECT SNOWFLAKE.CORTEX.COMPLETE(%s, %s)
    """, (cortex_model, full_prompt))

    result = cursor.fetchone()
    answer = result[0] if result else "I'm sorry, I couldn't generate a response."

    cursor.close()

    # Log AI answer
    log_to_snowflake(conv_id, user_id, True, answer, connection)

    # Return answer with citations
    response_data = {
        "answer": answer,
        "citations": citations
    }

    return True, "Successful message!", response_data


# --------------------------
# Generate Quiz Questions using Snowflake Cortex
# --------------------------
def generate_quiz(courseID, topic, difficulty, num_questions, connection, model=None, material_id=None):
    """
    Generate quiz questions using Snowflake Cortex AI based on course materials.

    Args:
        courseID: The course ID
        topic: Topic/chapter for the quiz
        difficulty: "Beginner", "Intermediate", or "Advanced"
        num_questions: Number of questions to generate (default 8)
        connection: Database connection
        model: Optional Cortex model override
        material_id: Optional specific material ID to generate quiz from

    Returns:
        (success, message, quiz_data)
        quiz_data format: {
            "title": "Quiz Title",
            "description": "Quiz Description",
            "questions": [
                {
                    "question": "Question text?",
                    "options": ["Option A", "Option B", "Option C", "Option D"],
                    "correct": 0,  # index of correct answer
                    "explanation": "Why this is correct"
                },
                ...
            ]
        }
    """
    cursor = connection.cursor()

    # If material_id is provided, fetch ONLY that specific material
    if material_id:
        print(f"DEBUG: Fetching material_id={material_id} for courseID={courseID}")
        cursor.execute("""
            SELECT title, content FROM course_materials
            WHERE material_id = %s AND course_id = %s
        """, (material_id, courseID))
        material = cursor.fetchone()

        if not material:
            cursor.close()
            print(f"DEBUG: Material {material_id} not found in database")
            return False, f"Material {material_id} not found", None

        print(f"DEBUG: Found material: {material[0]}, content length: {len(material[1]) if material[1] else 0}")

        title, content = material
        baseline_context = f"""You are Edwin, TA AI. Generate quiz questions based on this material:

MATERIAL: {title}
CONTENT: {content[:3000]}
"""
    else:
        # Use general course materials if no specific material
        baseline_context = get_baseline_context(courseID, connection, limit_chars=600)

    cursor.close()

    # Build quiz generation prompt (simplified for speed)
    prompt = f"""{baseline_context}

TASK: Generate {num_questions} {difficulty}-level multiple choice questions about: {topic}

REQUIREMENTS:
1. Each question must have exactly 4 options
2. Questions should test understanding, not just memorization
3. Include detailed explanations for correct answers
4. Base questions on the course materials provided above
5. Make questions progressively harder within the difficulty level

OUTPUT FORMAT (strict JSON only, no other text):
{{
    "title": "{topic}",
    "description": "Brief description of what this quiz covers",
    "questions": [
        {{
            "question": "Question text here?",
            "options": ["Option A", "Option B", "Option C", "Option D"],
            "correct": 0,
            "explanation": "Detailed explanation of why this answer is correct"
        }}
    ]
}}

Generate the quiz now:"""

    # Use Snowflake Cortex to generate quiz (use faster model)
    cortex_model = model or "mixtral-8x7b"  # Faster model for quiz generation
    cursor = connection.cursor()

    try:
        cursor.execute("""
            SELECT SNOWFLAKE.CORTEX.COMPLETE(%s, %s)
        """, (cortex_model, prompt))

        result = cursor.fetchone()
        response = result[0] if result else None

        if not response:
            cursor.close()
            return False, "Failed to generate quiz", None

        # Parse JSON response
        import json

        # Try to extract JSON from response (in case AI added extra text)
        if '{' in response:
            json_start = response.index('{')
            json_end = response.rindex('}') + 1
            json_str = response[json_start:json_end]
            quiz_data = json.loads(json_str)
        else:
            cursor.close()
            return False, "AI response was not valid JSON", None

        cursor.close()
        return True, "Quiz generated successfully", quiz_data

    except Exception as e:
        cursor.close()
        return False, f"Error generating quiz: {str(e)}", None


if __name__ == "__main__":
    from credentials import get_db_connection
    conn = get_db_connection()
    if conn:
        status, message, error = ingest_pdf_to_snowflake("Lab1CSE434.pdf", 231849, conn)
        status, message, error = ingest_pptx_to_snowflake("Chapter_2_v8.0-2.pptx", 231849, conn)
        log_to_snowflake("test_conv_123", 1, False, "Hello Edwin (and snowflake!)", conn)  # log user question
